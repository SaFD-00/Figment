"""ResolvedFigure → .pptx 바이트. python-pptx 1.0.2.

도형은 add_shape(매핑), 커넥터는 add_connector + begin/end_connect로 도형에 부착(드래그 시
화살표 추종). 모든 shape.name='fg-{id}'(3중 키). word_wrap=False/auto_size=NONE로 SVG와
줄바꿈 일치. 화살촉/점선/알파는 pptx_xml lxml 핵.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Mm, Pt

from ..units import BOX_ICON_MM
from . import pptx_xml
from .resolved import (
    RChart,
    RConnector,
    ResolvedFigure,
    RGroup,
    RImage,
    RShape,
    RText,
    TextRun,
)

SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "cylinder": MSO_SHAPE.CAN,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "hexagon": MSO_SHAPE.HEXAGON,
}
CONNECTOR_MAP = {
    "straight": MSO_CONNECTOR.STRAIGHT,
    "elbow": MSO_CONNECTOR.ELBOW,
    "curve": MSO_CONNECTOR.CURVE,
}
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
_CXN_IDX = {"top": 0, "left": 1, "bottom": 2, "right": 3}


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr.lstrip("#").upper())


class PptxRenderer:
    def __init__(self, asset_store: Any | None = None):
        self.asset_store = asset_store

    def render(self, fig: ResolvedFigure) -> bytes:
        prs = Presentation()
        prs.slide_width = Mm(fig.width_mm)
        prs.slide_height = Mm(fig.height_mm)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        shapes = slide.shapes

        # 배경 (맨 뒤)
        bg = shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(0), Mm(0), Mm(fig.width_mm), Mm(fig.height_mm))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(fig.background)
        bg.line.fill.background()
        pptx_xml.set_shape_name(bg, "fg-bg")

        id_to_shape: dict[str, Any] = {}
        connectors: list[RConnector] = []
        for el in fig.elements:
            if isinstance(el, RShape):
                id_to_shape[el.id] = self._shape(shapes, el)
            elif isinstance(el, RGroup):
                id_to_shape[el.id] = self._group(shapes, el)
            elif isinstance(el, RText):
                self._text(shapes, el)
            elif isinstance(el, RImage):
                self._image(shapes, el)
            elif isinstance(el, RChart):
                self._chart(shapes, el)
            elif isinstance(el, RConnector):
                connectors.append(el)

        for el in connectors:
            self._connector(shapes, el, id_to_shape)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── 도형 ──────────────────────────────────────────────────────────────────
    def _shape(self, shapes: Any, el: RShape) -> Any:
        shp = shapes.add_shape(
            SHAPE_MAP.get(el.shape_kind, MSO_SHAPE.RECTANGLE),
            Mm(el.x), Mm(el.y), Mm(el.w), Mm(el.h),
        )
        if el.fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(el.fill)
            pptx_xml.set_fill_alpha(shp, el.fill_opacity)
        else:
            shp.fill.background()
        if el.stroke_color:
            shp.line.color.rgb = _rgb(el.stroke_color)
            shp.line.width = Pt(el.stroke_width_pt)
            pptx_xml.set_dash(shp, el.dash)
        else:
            shp.line.fill.background()
        if el.shape_kind == "rounded":
            try:
                shp.adjustments[0] = max(0.0, min(0.5, el.corner_radius_mm / max(1.0, min(el.w, el.h))))
            except Exception:  # noqa: BLE001
                pass
        anchor = "middle"
        if getattr(el, "icon_asset", None) and self.asset_store is not None:  # M4.3
            png = _try_png(self.asset_store, el.icon_asset)
            if png:
                iw = min(BOX_ICON_MM, el.w - 2.0, max(4.0, el.h * 0.5))
                ix = el.x + (el.w - iw) / 2
                pic = shapes.add_picture(io.BytesIO(png), Mm(ix), Mm(el.y + 1.5), Mm(iw), Mm(iw))
                pptx_xml.set_shape_name(pic, f"fg-{el.id}__icon")
                anchor = "bottom"  # 라벨은 아이콘 아래
        self._fill_tf(shp.text_frame, el.label, anchor=anchor)
        pptx_xml.set_shape_name(shp, f"fg-{el.id}")
        return shp

    def _group(self, shapes: Any, el: RGroup) -> Any:
        shp = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Mm(el.x), Mm(el.y), Mm(el.w), Mm(el.h))
        shp.fill.background()
        shp.line.color.rgb = _rgb(el.stroke_color or "#9AA0A6")
        shp.line.width = Pt(el.stroke_width_pt)
        pptx_xml.set_dash(shp, el.dash)
        pptx_xml.set_shape_name(shp, f"fg-{el.id}")
        if el.label:
            tb = shapes.add_textbox(Mm(el.x + 2), Mm(el.y + 0.5), Mm(el.w - 4), Mm(el.label.line_height_mm + 1))
            self._fill_tf(tb.text_frame, el.label, anchor="top")
        return shp

    def _text(self, shapes: Any, el: RText) -> Any:
        tb = shapes.add_textbox(Mm(el.x), Mm(el.y), Mm(el.w), Mm(el.h))
        self._fill_tf(tb.text_frame, el.run, anchor=el.run.v_align)
        pptx_xml.set_shape_name(tb, f"fg-{el.id}")
        return tb

    def _image(self, shapes: Any, el: RImage) -> Any:
        png = None
        if el.asset_id and self.asset_store is not None:
            png = _try_png(self.asset_store, el.asset_id)
        if png:
            pic = shapes.add_picture(io.BytesIO(png), Mm(el.x), Mm(el.y), Mm(el.w), Mm(el.h))
            pptx_xml.set_shape_name(pic, f"fg-{el.id}")
            return pic
        return self._placeholder(shapes, el.id, el.x, el.y, el.w, el.h, el.placeholder_label)

    def _chart(self, shapes: Any, el: RChart) -> Any:
        png = None
        if el.svg_asset_id and self.asset_store is not None:
            png = _try_chart_png(self.asset_store, el.svg_asset_id)
        if png:
            pic = shapes.add_picture(io.BytesIO(png), Mm(el.x), Mm(el.y), Mm(el.w), Mm(el.h))
            pptx_xml.set_shape_name(pic, f"fg-{el.id}")
            return pic
        return self._placeholder(shapes, el.id, el.x, el.y, el.w, el.h, el.placeholder_label or "chart")

    def _placeholder(self, shapes, eid, x, y, w, h, label) -> Any:
        shp = shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(x), Mm(y), Mm(w), Mm(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb("#F1F3F4")
        shp.line.color.rgb = _rgb("#BDC1C6")
        shp.line.width = Pt(0.5)
        pptx_xml.set_dash(shp, "dash")
        run = TextRun(lines=[label], size_pt=7, color="#80868B", h_align="center",
                      v_align="middle", line_height_mm=3.0)
        shp.text_frame.word_wrap = True
        self._fill_tf(shp.text_frame, run, anchor="middle")
        pptx_xml.set_shape_name(shp, f"fg-{eid}")
        return shp

    def _connector(self, shapes: Any, el: RConnector, id_to_shape: dict[str, Any]) -> Any:
        p0, p1 = el.points[0], el.points[-1]
        conn = shapes.add_connector(
            CONNECTOR_MAP.get(el.routing, MSO_CONNECTOR.STRAIGHT),
            Mm(p0[0]), Mm(p0[1]), Mm(p1[0]), Mm(p1[1]),
        )
        conn.line.color.rgb = _rgb(el.stroke_color)
        conn.line.width = Pt(el.stroke_width_pt)
        pptx_xml.set_dash(conn, el.dash)
        pptx_xml.set_line_arrowheads(conn, el.head, el.tail)
        pptx_xml.set_shape_name(conn, f"fg-{el.id}")
        # 도형 부착 (드래그 추종)
        try:
            if el.from_id in id_to_shape:
                conn.begin_connect(id_to_shape[el.from_id], _CXN_IDX.get(el.src_side or "right", 3))
            if el.to_id in id_to_shape:
                conn.end_connect(id_to_shape[el.to_id], _CXN_IDX.get(el.tgt_side or "left", 1))
        except Exception:  # noqa: BLE001
            pass
        if el.label and el.label_anchor:
            ax, ay = el.label_anchor
            tb = shapes.add_textbox(Mm(ax - 12), Mm(ay - 2), Mm(24), Mm(4))
            self._fill_tf(tb.text_frame, el.label, anchor="middle")
        return conn

    # ── 텍스트 프레임 ──────────────────────────────────────────────────────────
    def _fill_tf(self, tf: Any, run: TextRun | None, *, anchor: str = "middle") -> None:
        tf.word_wrap = False
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:  # noqa: BLE001
            pass
        tf.margin_left = tf.margin_right = Mm(1)
        tf.margin_top = tf.margin_bottom = Mm(0.5)
        tf.vertical_anchor = _ANCHOR.get(anchor, MSO_ANCHOR.MIDDLE)
        if run is None or not run.lines:
            return
        for i, line in enumerate(run.lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = _ALIGN.get(run.h_align, PP_ALIGN.CENTER)
            r = p.add_run()
            r.text = line
            f = r.font
            f.size = Pt(run.size_pt)
            f.name = run.family
            f.bold = run.weight == "bold"
            f.italic = run.italic
            f.color.rgb = _rgb(run.color)


def _try_png(store: Any, asset_id: str) -> bytes | None:
    fn = getattr(store, "get_png", None)
    if fn:
        try:
            return fn(asset_id)
        except Exception:  # noqa: BLE001
            return None
    return None


def _try_chart_png(store: Any, asset_id: str) -> bytes | None:
    for meth in ("get_chart_png", "get_png"):
        fn = getattr(store, meth, None)
        if fn:
            try:
                return fn(asset_id)
            except Exception:  # noqa: BLE001
                continue
    return None
