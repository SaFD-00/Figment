"""렌더 동등성 — 골든 spec → SVG/PPTX 구조 스냅샷 + 결정성.

플랫폼 폰트 렌더 차이로 픽셀 diff는 불안정하므로, (1) 결정성(동일 spec→동일 SVG)과
(2) 구조 계약(data-fg-id / shape.name 3중 키, EMU 좌표, 화살촉, 텍스트 벡터)을 검증한다.
"""

from __future__ import annotations

import io
import re
from pathlib import Path

import pytest
from pptx import Presentation

from figgen.layout import LayoutEngine
from figgen.render.pptx_renderer import PptxRenderer
from figgen.render.resolver import resolve
from figgen.render.svg_renderer import SvgRenderer
from figgen.schema import FigureSpec
from figgen.styles.presets import get_preset

GOLDEN = Path(__file__).parent / "fixtures" / "golden" / "method_basic.spec.json"
_VISUAL = ("box", "text", "image", "chart", "group")


@pytest.fixture(scope="module")
def fig():
    spec = FigureSpec.model_validate_json(GOLDEN.read_text(encoding="utf-8"))
    spec = spec.model_copy(update={"stylesheet": get_preset("nature_minimal")})
    layout = LayoutEngine().layout(spec)
    return spec, resolve(spec, layout, spec.stylesheet)


def _visual_ids(spec: FigureSpec) -> set[str]:
    """렌더되는(선택 가능한) 요소 id — 비시각 컨테이너(row/column/grid) 제외 + 커넥터."""
    ids = {n.id for n, _ in spec.iter_elements() if getattr(n, "type", None) in _VISUAL}
    return ids | {c.id for c in spec.connectors}


def test_svg_has_all_fg_ids(fig):
    spec, f = fig
    svg = SvgRenderer().render(f)
    ids = set(re.findall(r'data-fg-id="([^"]+)"', svg))
    # 모든 요소(컨테이너 중 row는 비시각 → 제외)와 커넥터가 등장
    expected = {"encoder", "decoder", "input", "embed", "attn", "ffn", "cross", "dffn", "output"}
    expected |= {"c_in", "c_enc_dec", "c_out", "c_fb"}
    assert expected <= ids


def test_svg_deterministic(fig):
    _, f = fig
    a = SvgRenderer().render(f)
    b = SvgRenderer().render(f)
    assert a == b  # 바이트 동일


def test_svg_text_is_vector(fig):
    _, f = fig
    svg = SvgRenderer().render(f)
    assert "<text" in svg and "<tspan" in svg
    assert "Multi-Head Self-Attention" in svg  # 라벨이 텍스트로 존재
    assert 'viewBox="0 0 180' in svg


def test_svg_has_arrow_markers(fig):
    _, f = fig
    svg = SvgRenderer().render(f)
    assert "marker-end" in svg and "<marker" in svg


def test_pptx_shape_names_match_ids(fig):
    spec, f = fig
    data = PptxRenderer().render(f)
    prs = Presentation(io.BytesIO(data))
    names = {sh.name for sh in prs.slides[0].shapes}
    for eid in _visual_ids(spec):
        assert f"fg-{eid}" in names, f"fg-{eid} 누락"


def test_pptx_slide_size_emu(fig):
    _, f = fig
    prs = Presentation(io.BytesIO(PptxRenderer().render(f)))
    # 180mm = 6,480,000 EMU
    assert prs.slide_width == 6_480_000
    assert abs(prs.slide_height - f.height_mm * 36000) < 2


def test_pptx_has_arrowheads_and_connectors(fig):
    _, f = fig
    prs = Presentation(io.BytesIO(PptxRenderer().render(f)))
    xml = prs.slides[0].shapes._spTree.xml
    assert "tailEnd" in xml  # 화살촉
    assert "cxnSp" in xml  # 커넥터
    assert "prstDash" in xml  # 점선(feedback)


def test_render_does_not_mutate_spec(fig):
    spec, f = fig
    before = spec.model_dump_json()
    SvgRenderer().render(f)
    PptxRenderer().render(f)
    assert spec.model_dump_json() == before
