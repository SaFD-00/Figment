"""Export / vectorize helpers — turn any raster asset into editable formats.

  • png_to_svg  → vtracer multicolor path SVG (via the vendored FigGen vectorizer).
  • png_to_pptx → a single full-bleed slide embedding the image (python-pptx).

Figure-engine assets already carry editable figure.svg / figure.pptx sidecars (in asset.meta);
these helpers cover plain raster outputs (ComfyUI / cloud images) so every asset is exportable.
"""
from __future__ import annotations

import io

from figgen.fullimage.vectorize import vectorize_png

PPTX_MEDIA = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
SVG_MEDIA = "image/svg+xml"
EMU_PER_INCH = 914400
LONG_SIDE_IN = 10.0  # scale the image's longer side to 10 inches on the slide


def png_to_svg(png: bytes) -> str:
    """Deterministic raster → layered multicolor SVG string."""
    return vectorize_png(png)


def png_to_pptx(png: bytes) -> bytes:
    """Embed a PNG as a single full-slide image in a .pptx, sized to the image aspect.

    The slide is scaled so the longer side is 10 inches (python-pptx requires slide
    dimensions within 1–56 inches), preserving the image aspect ratio.
    """
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image

    with Image.open(io.BytesIO(png)) as im:
        w, h = im.size
    w = max(w, 1)
    h = max(h, 1)

    if w >= h:
        sw_in, sh_in = LONG_SIDE_IN, LONG_SIDE_IN * h / w
    else:
        sh_in, sw_in = LONG_SIDE_IN, LONG_SIDE_IN * w / h
    sw_in = max(sw_in, 1.0)
    sh_in = max(sh_in, 1.0)

    prs = Presentation()
    prs.slide_width = Emu(int(sw_in * EMU_PER_INCH))
    prs.slide_height = Emu(int(sh_in * EMU_PER_INCH))
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.shapes.add_picture(io.BytesIO(png), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
